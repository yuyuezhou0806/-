from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

IMG_DIR = r"c:\Users\admin\cc_test"

def add_top_bar(slide, color=RGBColor(0x1a, 0x5f, 0x7a)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

def add_title(slide, text, top=Inches(0.25), color=RGBColor(0x1a, 0x5f, 0x7a)):
    tb = slide.shapes.add_textbox(Inches(0.4), top, Inches(12.5), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = color
    return tb

def add_bullet_text(slide, items, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_before = Pt(12)
        p.space_after = Pt(6)
    return tb

def add_image_slide(slide, img_path, img_left, img_top, max_width, max_height):
    """Add image with aspect ratio preserved"""
    try:
        with Image.open(img_path) as img:
            orig_w, orig_h = img.size
        aspect = orig_w / orig_h

        if max_width / max_height > aspect:
            h = max_height
            w = h * aspect
        else:
            w = max_width
            h = w / aspect

        pic = slide.shapes.add_picture(img_path, img_left, img_top, width=w, height=h)
        return pic
    except Exception as e:
        print(f"Error adding image {img_path}: {e}")
        return None

# ========== 第1页：封面 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = RGBColor(0x1a, 0x5f, 0x7a)
bg.line.fill.background()
tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.333), Inches(1.5))
p = tb.text_frame.paragraphs[0]
p.text = "客户申领流程"
p.font.size = Pt(54)
p.font.bold = True
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p.alignment = PP_ALIGN.CENTER
p2 = tb.text_frame.add_paragraph()
p2.text = "从客户公海到提交审批"
p2.font.size = Pt(28)
p2.font.color.rgb = RGBColor(0xCC, 0xE5, 0xFF)
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(20)

# ========== 第2页：进入营销管理 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_top_bar(slide)
add_title(slide, "第一步：进入营销管理")
add_bullet_text(slide, [
    "1. 点击顶部导航栏「更多应用」",
    "2. 在弹出菜单中找到「营销管理」",
    "3. 点击进入"
], Inches(0.4), Inches(1.1), Inches(5.5), Inches(3))
add_image_slide(slide, f"{IMG_DIR}\\eaf2ec36fe7cc31f0a997c9007560365.png",
                Inches(6.2), Inches(1.0), Inches(6.8), Inches(6.2))

# ========== 第3页：进入客户公海 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_top_bar(slide)
add_title(slide, "第二步：进入客户公海")
add_bullet_text(slide, [
    "1. 左侧菜单展开「客户」",
    "2. 点击「客户公海」",
    "3. 进入公海客户列表"
], Inches(0.4), Inches(1.1), Inches(5.5), Inches(3))
add_image_slide(slide, f"{IMG_DIR}\\2caae3fd2222dfa0586e460354bf3884.png",
                Inches(6.2), Inches(1.0), Inches(6.8), Inches(6.2))

# ========== 第4页：搜索客户 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_top_bar(slide)
add_title(slide, "第三步：搜索目标客户")
add_bullet_text(slide, [
    "1. 在顶部搜索框输入客户名称",
    "   例：上海市道路运输事业发展中心",
    "2. 点击右侧「搜索」按钮",
    "3. 查看搜索结果"
], Inches(0.4), Inches(1.1), Inches(5.5), Inches(3.5))
add_image_slide(slide, f"{IMG_DIR}\\80cbf3f71bc5c90431e5b179e6fe6442.png",
                Inches(6.2), Inches(1.0), Inches(6.8), Inches(6.2))

# ========== 第5页：发起客户申领 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_top_bar(slide)
add_title(slide, "第四步：发起客户申领")
add_bullet_text(slide, [
    "1. 在搜索结果中找到目标客户",
    "2. 点击「客户申领」按钮",
    "3. 系统弹出申领流程表单"
], Inches(0.4), Inches(1.1), Inches(5.5), Inches(3))
add_image_slide(slide, f"{IMG_DIR}\\a8f36e10a2eff18b6f6702f3b8558d58.png",
                Inches(6.2), Inches(1.0), Inches(6.8), Inches(6.2))

# ========== 第6页：填写申领信息 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_top_bar(slide)
add_title(slide, "第五步：填写申领信息")
add_bullet_text(slide, [
    "1. 客户名称：自动带入",
    "2. 客户经理：点击选择负责人",
    "   例：卫天勇",
    "3. 人员、部门自动带出",
    "4. 确认新领取人员"
], Inches(0.4), Inches(1.1), Inches(5.5), Inches(3.5))
add_image_slide(slide, f"{IMG_DIR}\\a77456aa77e80f2bdad66a6a1345e154.png",
                Inches(6.2), Inches(1.0), Inches(6.8), Inches(6.2))

# ========== 第7页：提交审批 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_top_bar(slide)
add_title(slide, "第六步：提交审批")
add_bullet_text(slide, [
    "1. 确认所有信息无误",
    "2. 点击右上角「提交」按钮",
    "3. 等待审批通过",
    "4. 通过后客户进入「我的客户」"
], Inches(0.4), Inches(1.1), Inches(5.5), Inches(3.5))
add_image_slide(slide, f"{IMG_DIR}\\a77456aa77e80f2bdad66a6a1345e154.png",
                Inches(6.2), Inches(1.0), Inches(6.8), Inches(6.2))

# ========== 第8页：流程总结 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_top_bar(slide, RGBColor(0x28, 0xA7, 0x45))
tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.6))
p = tb.text_frame.paragraphs[0]
p.text = "流程总结"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = RGBColor(0x28, 0xA7, 0x45)

add_bullet_text(slide, [
    "更多应用 → 营销管理 → 客户 → 客户公海",
    "→ 输入客户名称 → 搜索",
    "→ 客户申领 → 选择客户经理 → 提交",
    "",
    "审批通过后客户即转入「我的客户」"
], Inches(0.4), Inches(1.2), Inches(12.5), Inches(4))

# 保存
output = r"c:\Users\admin\cc_test\客户申领流程.pptx"
prs.save(output)
print(f"PPT已生成：{output}")
