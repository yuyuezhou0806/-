from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_slide(prs, title, items, title_color=RGBColor(0x1a, 0x5f, 0x7a)):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = title_color
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = title_color

    cb = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12.333), Inches(5.5))
    tf = cb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(26)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_before = Pt(20)
        p.space_after = Pt(10)
    return slide

def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x1a, 0x5f, 0x7a)
    bg.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        p2 = tb.text_frame.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(28)
        p2.font.color.rgb = RGBColor(0xCC, 0xE5, 0xFF)
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(20)
    return slide

# 1. 封面
add_title_slide(prs, "客户申领流程", "从客户公海到提交审批")

# 2. 入口
add_slide(prs, "第一步：进入营销管理", [
    "1. 点击顶部导航栏「更多应用」",
    "2. 在弹出菜单中找到「营销管理」模块",
    "3. 点击进入"
])

# 3. 进入客户公海
add_slide(prs, "第二步：进入客户公海", [
    "1. 左侧菜单展开「客户」",
    "2. 点击「客户公海」",
    "3. 进入公海客户列表页面"
])

# 4. 搜索客户
add_slide(prs, "第三步：搜索目标客户", [
    "1. 在顶部搜索框输入客户名称",
    "   例：上海市道路运输事业发展中心",
    "2. 点击右侧「搜索」按钮",
    "3. 查看搜索结果"
])

# 5. 发起客户申领
add_slide(prs, "第四步：发起客户申领", [
    "1. 在搜索结果中找到目标客户",
    "2. 点击「客户申领」按钮",
    "3. 系统弹出申领流程表单"
])

# 6. 填写申领信息
add_slide(prs, "第五步：填写申领信息", [
    "1. 客户名称：自动带入搜索的客户",
    "2. 客户经理：点击选择框，选择负责人",
    "   例：卫天勇",
    "3. 人员：选择相关人员",
    "4. 部门：自动带出客户经理所属部门",
    "   例：中测行运营室",
    "5. 新领取人员：确认或调整"
])

# 7. 提交
add_slide(prs, "第六步：提交审批", [
    "1. 确认信息无误",
    "2. 点击右上角「提交」按钮",
    "3. 等待审批通过后，该客户进入「我的客户」"
])

# 8. 流程总结
add_slide(prs, "流程总结", [
    "更多应用 → 营销管理 → 客户 → 客户公海",
    "→ 输入客户名称 → 搜索",
    "→ 客户申领 → 选择客户经理 → 提交",
    "",
    "审批通过后客户即转入个人客户列表"
], RGBColor(0x28, 0xA7, 0x45))

prs.save(r"c:\Users\admin\cc_test\客户申领流程.pptx")
print("PPT已生成：c:\\Users\\admin\\cc_test\\客户申领流程.pptx")
