"""
新平台开发需求 PPT 生成脚本
基于 C:\\Users\\admin\\Desktop\\scan\\新平台开发需求.docx 内容生成
目标:与软件公司沟通使用
"""
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy

IMG_DIR = r"C:\Users\admin\cc_test\imgs"
def img(name):
    return os.path.join(IMG_DIR, name)

# ---------- 配色 ----------
PRIMARY = RGBColor(0x1F, 0x3A, 0x68)      # 深蓝(标题主色)
ACCENT = RGBColor(0x2E, 0x86, 0xC1)        # 亮蓝(强调)
WARN = RGBColor(0xC0, 0x39, 0x2B)          # 红色(痛点)
SUCCESS = RGBColor(0x1E, 0x82, 0x49)       # 绿色(目标/收益)
GREY_DARK = RGBColor(0x33, 0x33, 0x33)
GREY_MED = RGBColor(0x66, 0x66, 0x66)
GREY_LIGHT = RGBColor(0xF2, 0xF4, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT_TITLE = "微软雅黑"
FONT_BODY = "微软雅黑"

# ---------- 工具函数 ----------
def set_run_font(run, name=FONT_BODY, size=18, bold=False, color=GREY_DARK):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    # 中文字体需要单独设置 eastAsia
    rPr = run._r.get_or_add_rPr()
    rFonts = rPr.find(qn('a:rFonts'))
    if rFonts is None:
        from lxml import etree
        rFonts = etree.SubElement(rPr, qn('a:rFonts'))
    rFonts.set(qn('a:eastAsia'), name)

def add_textbox(slide, left, top, width, height, text, size=18, bold=False,
                color=GREY_DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT_BODY):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run_font(run, font, size, bold, color)
    return tb

def add_bullets(slide, left, top, width, height, items, size=16, color=GREY_DARK,
                line_spacing=1.3, bullet_color=ACCENT):
    """items: list[str] or list[(title, desc)]"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_top = Emu(0)
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        # bullet 符号
        run0 = p.add_run()
        run0.text = "■ "
        set_run_font(run0, FONT_BODY, size, True, bullet_color)
        if isinstance(item, tuple):
            title, desc = item
            r1 = p.add_run()
            r1.text = title
            set_run_font(r1, FONT_BODY, size, True, GREY_DARK)
            if desc:
                r2 = p.add_run()
                r2.text = "  " + desc
                set_run_font(r2, FONT_BODY, size - 1, False, GREY_MED)
        else:
            r1 = p.add_run()
            r1.text = item
            set_run_font(r1, FONT_BODY, size, False, color)
    return tb

def add_rect(slide, left, top, width, height, fill=GREY_LIGHT, line=None, line_w=0):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w)
    shape.shadow.inherit = False
    return shape

def add_rounded(slide, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    shape.shadow.inherit = False
    return shape

def add_card(slide, left, top, width, height, title, body, title_color=PRIMARY,
             body_size=13, title_size=18):
    add_rounded(slide, left, top, width, height, WHITE, line=ACCENT)
    add_textbox(slide, left + Inches(0.15), top + Inches(0.1),
                width - Inches(0.3), Inches(0.45), title, size=title_size, bold=True,
                color=title_color)
    add_textbox(slide, left + Inches(0.15), top + Inches(0.6),
                width - Inches(0.3), height - Inches(0.7), body, size=body_size,
                color=GREY_DARK)

def fit_picture(slide, path, left, top, max_w, max_h):
    """按图片原始比例缩放到 box 内,并居中放置。max_w/max_h 为 Emu。"""
    with Image.open(path) as im:
        iw, ih = im.size
    ratio = iw / ih
    target_w = max_w
    target_h = int(target_w / ratio)
    if target_h > max_h:
        target_h = max_h
        target_w = int(target_h * ratio)
    cx = left + (max_w - target_w) // 2
    cy = top + (max_h - target_h) // 2
    pic = slide.shapes.add_picture(path, cx, cy, width=target_w, height=target_h)
    return pic

def add_module_block(slide, x, y, w, h, title, image_files, label_color=PRIMARY):
    """在矩形区域内,顶部放模块标题,下方按行平铺图片(自动按图数量调整高度)。"""
    # 标题条
    add_rect(slide, x, y, w, Inches(0.4), fill=label_color)
    add_textbox(slide, x + Inches(0.15), y + Inches(0.02), w - Inches(0.3),
                Inches(0.36), title, size=14, bold=True, color=WHITE,
                anchor=MSO_ANCHOR.MIDDLE)
    # 图区
    img_top = y + Inches(0.5)
    img_area_h = h - Inches(0.5)
    n = len(image_files)
    if n == 0:
        return
    gap = Inches(0.08)
    each_h = (img_area_h - gap * (n - 1)) // n
    for i, f in enumerate(image_files):
        iy = img_top + i * (each_h + gap)
        # 边框
        add_rect(slide, x, iy, w, each_h, fill=WHITE, line=GREY_MED, line_w=0.5)
        fit_picture(slide, img(f), x + Inches(0.05), iy + Inches(0.05),
                    w - Inches(0.1), each_h - Inches(0.1))

def add_page_title(slide, title, subtitle=None):
    # 顶部装饰条
    bar = add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), fill=PRIMARY)
    add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
                title, size=26, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(1.0), Inches(12), Inches(0.4),
                    subtitle, size=14, color=GREY_MED)

def add_footer(slide, page_no, total):
    add_rect(slide, Inches(0), Inches(7.3), Inches(13.333), Inches(0.2), fill=PRIMARY)
    add_textbox(slide, Inches(0.5), Inches(7.05), Inches(8), Inches(0.25),
                "新平台开发需求 · 内部沟通材料", size=10, color=GREY_MED)
    add_textbox(slide, Inches(11.5), Inches(7.05), Inches(1.5), Inches(0.25),
                f"{page_no} / {total}", size=10, color=GREY_MED, align=PP_ALIGN.RIGHT)

# ---------- 创建演示文稿 ----------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

TOTAL = 19

# ===================== 1. 封面 =====================
s = prs.slides.add_slide(blank)
add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5), fill=PRIMARY)
# 装饰
add_rect(s, Inches(0), Inches(5.5), Inches(13.333), Inches(0.05), fill=ACCENT)
add_textbox(s, Inches(1), Inches(2.3), Inches(11.3), Inches(1),
            "新合同管理平台开发需求", size=44, bold=True, color=WHITE,
            font=FONT_TITLE)
add_textbox(s, Inches(1), Inches(3.3), Inches(11.3), Inches(0.7),
            "Contract Management Platform — Requirement Proposal",
            size=18, color=RGBColor(0xC9, 0xD6, 0xE6))
add_textbox(s, Inches(1), Inches(5.7), Inches(11.3), Inches(0.5),
            "用于与软件开发公司前期沟通", size=16, color=WHITE)
add_textbox(s, Inches(1), Inches(6.3), Inches(11.3), Inches(0.4),
            "2026 / 05", size=14, color=RGBColor(0xC9, 0xD6, 0xE6))

# ===================== 2. 议程 =====================
s = prs.slides.add_slide(blank)
add_page_title(s, "议程  Agenda")
items = [
    ("01  业务背景", "当前合同从立项到归档的整体流程"),
    ("02  现状痛点", "六个平台重复录入,管理平台五大局限"),
    ("03  新平台目标", "替代内网管理平台,自动对接外网协会与 BPM"),
    ("04  核心功能", "OCR 抓取 / 合同生成 / 电子流转 / BPM 同步 / AI 分析"),
    ("05  新旧流程对比", "老方法流程图 vs 新平台流程图"),
    ("06  网络拓扑", "内外网部署、跨网数据流转方案"),
    ("07  待讨论 + 后续", "技术路线 / 对接方式 / 实施节奏"),
]
y = Inches(1.4)
for t, d in items:
    add_rect(s, Inches(0.8), y, Inches(0.1), Inches(0.65), fill=ACCENT)
    add_textbox(s, Inches(1.1), y + Inches(0.02), Inches(3.5), Inches(0.4),
                t, size=19, bold=True, color=PRIMARY)
    add_textbox(s, Inches(1.1), y + Inches(0.38), Inches(11), Inches(0.3),
                d, size=12, color=GREY_MED)
    y += Inches(0.78)
add_footer(s, 2, TOTAL)

# ===================== 3. 业务背景:同一份合同被录入 6 处 =====================
s = prs.slides.add_slide(blank)
add_page_title(s, "一、业务背景", "一份合同信息,需要在 6 个以上平台 / 文档中重复录入")

# 中心源
center_x, center_y = Inches(5.8), Inches(3.5)
src_w, src_h = Inches(2.2), Inches(1.3)
add_rounded(s, center_x, center_y, src_w, src_h, WARN)
add_textbox(s, center_x, center_y + Inches(0.2), src_w, Inches(0.4),
            "信息源", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(s, center_x, center_y + Inches(0.6), src_w, Inches(0.5),
            "施工许可证", size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# 6 个目标平台
targets = [
    ("行业协会平台", "录入 + 勾选参数", "外网 · 老旧"),
    ("合同文本制作", "Word 模板填写", "本地"),
    ("管理平台", "录入 + 合同归档", "内网 · 老旧"),
    ("BPM 商机", "商机数据录入", "内/外网均可登录"),
    ("BPM 销售合同", "销售合同录入", "内/外网均可登录"),
    ("合同流转单", "流转单填写", "内部"),
]
positions = [
    (Inches(0.5), Inches(1.5)),
    (Inches(0.5), Inches(3.4)),
    (Inches(0.5), Inches(5.3)),
    (Inches(10.6), Inches(1.5)),
    (Inches(10.6), Inches(3.4)),
    (Inches(10.6), Inches(5.3)),
]
for (px, py), (name, desc, net) in zip(positions, targets):
    # 管理平台用红框强调(替代目标)
    is_target = (name == "管理平台")
    border = WARN if is_target else ACCENT
    add_rounded(s, px, py, Inches(2.2), Inches(1.5), WHITE, line=border)
    add_textbox(s, px, py + Inches(0.15), Inches(2.2), Inches(0.4),
                name, size=14, bold=True,
                color=WARN if is_target else PRIMARY, align=PP_ALIGN.CENTER)
    add_textbox(s, px, py + Inches(0.55), Inches(2.2), Inches(0.4),
                desc, size=11, color=GREY_MED, align=PP_ALIGN.CENTER)
    # 网络标签
    tag_w = Inches(1.6)
    tag_x = px + (Inches(2.2) - tag_w) / 2
    add_rounded(s, tag_x, py + Inches(1.0), tag_w, Inches(0.32),
                GREY_LIGHT)
    add_textbox(s, tag_x, py + Inches(1.02), tag_w, Inches(0.3),
                net, size=10, bold=True, color=PRIMARY,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 连线
    line = s.shapes.add_connector(1,
        center_x + src_w / 2, center_y + src_h / 2,
        px + Inches(1.1), py + Inches(0.65))
    line.line.color.rgb = ACCENT
    line.line.width = Pt(1.2)

# 底部小结
add_rect(s, Inches(2.0), Inches(6.95), Inches(9.3), Inches(0.4), fill=GREY_LIGHT)
add_textbox(s, Inches(2.0), Inches(6.97), Inches(9.3), Inches(0.35),
            "  同一组核心信息反复录入 6 处,跨越内网管理平台 / 外网协会网 / BPM 多模块 → 新平台需在内网替代管理平台并打通各系统",
            size=12, bold=True, color=WARN, anchor=MSO_ANCHOR.MIDDLE)
add_footer(s, 3, TOTAL)

# ===================== 4. 施工许可证字段(信息源拆解) =====================
s = prs.slides.add_slide(blank)
add_page_title(s, "二、信息源拆解 · 施工许可证", "可被自动抓取的关键字段")

left_items = [
    ("报建号", "工程立项编号,合同主键"),
    ("施工许可证编号", "证件编号,流程查验依据"),
    ("工程名称", "合同名称基础"),
    ("建设地址", "项目地点"),
]
right_items = [
    ("建设单位", "甲方/委托方"),
    ("施工单位", "总包"),
    ("设计单位", "设计方"),
    ("监理单位", "监理方"),
]

def field_card(slide, left, top, name, desc):
    add_rounded(slide, left, top, Inches(5.5), Inches(0.95), WHITE, line=ACCENT)
    add_rect(slide, left, top, Inches(0.1), Inches(0.95), fill=ACCENT)
    add_textbox(slide, left + Inches(0.25), top + Inches(0.12),
                Inches(5.2), Inches(0.4), name, size=16, bold=True, color=PRIMARY)
    add_textbox(slide, left + Inches(0.25), top + Inches(0.5),
                Inches(5.2), Inches(0.4), desc, size=12, color=GREY_MED)

y = Inches(1.5)
for (n, d) in left_items:
    field_card(s, Inches(0.7), y, n, d)
    y += Inches(1.1)
y = Inches(1.5)
for (n, d) in right_items:
    field_card(s, Inches(7.1), y, n, d)
    y += Inches(1.1)

add_rect(s, Inches(0.7), Inches(6.2), Inches(11.9), Inches(0.7), fill=GREY_LIGHT)
add_textbox(s, Inches(0.9), Inches(6.3), Inches(11.6), Inches(0.5),
            "需 OCR / 文档解析支持三种输入: 施工许可证图片  ·  住建委系统截图  ·  施工许可证 PDF",
            size=14, bold=True, color=PRIMARY, anchor=MSO_ANCHOR.MIDDLE)
add_footer(s, 4, TOTAL)

# ===================== 5. 老管理平台五大痛点 =====================
s = prs.slides.add_slide(blank)
add_page_title(s, "三、现有管理平台痛点", "内网 · 老旧系统,新平台的核心目标就是替代它")

pains = [
    ("01", "无法二次开发", "封闭系统,业务流程变化无法跟进"),
    ("02", "无法运行自动化脚本", "重复工作只能人工执行"),
    ("03", "数据采集不完整", "字段口径不齐,后续分析难做"),
    ("04", "历史参数堆积", "系统承载量过大,响应慢"),
    ("05", "无法跨平台自动联动", "所有信息靠手工填入 — 最痛"),
]
y = Inches(1.5)
for i, (no, t, d) in enumerate(pains):
    is_critical = (i == 4)
    bg = WARN if is_critical else WHITE
    fg_t = WHITE if is_critical else PRIMARY
    fg_d = RGBColor(0xFF, 0xE0, 0xDC) if is_critical else GREY_MED
    add_rounded(s, Inches(0.7), y, Inches(11.9), Inches(0.9), bg, line=WARN)
    add_textbox(s, Inches(0.9), y + Inches(0.15), Inches(0.9), Inches(0.6),
                no, size=24, bold=True,
                color=WARN if not is_critical else WHITE)
    add_textbox(s, Inches(1.9), y + Inches(0.1), Inches(4), Inches(0.5),
                t, size=18, bold=True, color=fg_t)
    add_textbox(s, Inches(6.2), y + Inches(0.18), Inches(6.4), Inches(0.5),
                d, size=14, color=fg_d)
    y += Inches(1.05)
add_footer(s, 5, TOTAL)

# ===================== 6. 新平台目标 =====================
s = prs.slides.add_slide(blank)
add_page_title(s, "四、新平台目标", "替代内网老管理平台,自动对接外网协会网与 BPM")

goals = [
    ("替代管理平台", "新平台部署于内网,完整接管现有合同管理职能"),
    ("一次录入", "施工许可证 OCR 识别 + 人工补齐价格 / 折扣 / 配合费"),
    ("自动同步", "向外网协会网备案、BPM 商机 / 销售合同推送 — 由平台自动完成"),
    ("电子流转", "审批流 + 电子签证,替代纸质流转单"),
    ("可二次开发", "开放 API + 脚本能力,业务变化能跟上"),
    ("数据闭环 + AI", "字段口径统一沉淀,每月自动数据分析报告"),
]
cols = 3
card_w = Inches(4.0)
card_h = Inches(2.4)
gap = Inches(0.2)
start_x = Inches(0.6)
start_y = Inches(1.5)
for i, (t, d) in enumerate(goals):
    r, c = divmod(i, cols)
    x = start_x + c * (card_w + gap)
    y = start_y + r * (card_h + gap)
    add_rounded(s, x, y, card_w, card_h, WHITE, line=SUCCESS)
    add_rect(s, x, y, card_w, Inches(0.08), fill=SUCCESS)
    add_textbox(s, x + Inches(0.3), y + Inches(0.25), card_w - Inches(0.6),
                Inches(0.6), t, size=22, bold=True, color=SUCCESS)
    add_textbox(s, x + Inches(0.3), y + Inches(0.95), card_w - Inches(0.6),
                card_h - Inches(1.1), d, size=13, color=GREY_DARK)
add_footer(s, 6, TOTAL)

# ===================== 7. 核心功能 1 — OCR 识别 =====================
s = prs.slides.add_slide(blank)
add_page_title(s, "五 · 核心功能 ① · 施工许可证识别", "支持图片 / 截图 / PDF 多种输入")

# 左侧:输入
add_textbox(s, Inches(0.7), Inches(1.4), Inches(4), Inches(0.5),
            "输入形式", size=18, bold=True, color=PRIMARY)
inputs = ["施工许可证 · 实物拍照", "住建委系统 · 网页截图", "施工许可证 · PDF 文件"]
y = Inches(2.0)
for t in inputs:
    add_rounded(s, Inches(0.7), y, Inches(4), Inches(0.7), GREY_LIGHT)
    add_textbox(s, Inches(0.9), y + Inches(0.18), Inches(3.6), Inches(0.4),
                t, size=14, bold=True, color=GREY_DARK)
    y += Inches(0.85)

# 箭头
arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(5.0), Inches(3.6),
                         Inches(1.0), Inches(0.6))
arr.fill.solid()
arr.fill.fore_color.rgb = ACCENT
arr.line.fill.background()

# 右侧:输出字段
add_textbox(s, Inches(6.4), Inches(1.4), Inches(6), Inches(0.5),
            "自动抓取字段", size=18, bold=True, color=PRIMARY)
fields = ["报建号", "工程名称", "建设单位", "建设地址",
          "设计单位", "施工单位", "监理单位", "施工许可证编号"]
y = Inches(2.0)
for i, f in enumerate(fields):
    r, c = divmod(i, 2)
    add_rounded(s, Inches(6.4) + c * Inches(3.1), y + r * Inches(0.7),
                Inches(2.9), Inches(0.55), WHITE, line=ACCENT)
    add_textbox(s, Inches(6.4) + c * Inches(3.1),
                y + r * Inches(0.7) + Inches(0.13),
                Inches(2.9), Inches(0.4),
                f, size=13, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)

# 底部备注
add_rect(s, Inches(0.7), Inches(6.4), Inches(11.9), Inches(0.55), fill=GREY_LIGHT)
add_textbox(s, Inches(0.9), Inches(6.45), Inches(11.6), Inches(0.45),
            "技术要求:识别准确率需保证可用,需提供人工校对界面;对扫描质量差的样本需有容错。",
            size=13, color=GREY_DARK, anchor=MSO_ANCHOR.MIDDLE)
add_footer(s, 7, TOTAL)

# ===================== 8. 核心功能 2 + 3 合同生成 / 电子流转 =====================
s = prs.slides.add_slide(blank)
add_page_title(s, "五 · 核心功能 ② ③ · 合同生成与电子流转")

# 左卡片
add_rounded(s, Inches(0.7), Inches(1.5), Inches(5.8), Inches(5.3), WHITE, line=ACCENT)
add_rect(s, Inches(0.7), Inches(1.5), Inches(5.8), Inches(0.6), fill=ACCENT)
add_textbox(s, Inches(0.9), Inches(1.55), Inches(5.4), Inches(0.5),
            "② 合同与备案表自动生成", size=18, bold=True, color=WHITE,
            anchor=MSO_ANCHOR.MIDDLE)
left_bullets = [
    "OCR 字段 + 人工补充(价格 / 折扣 / 配合费 / 收费方式)",
    "按合同类型选用 Word 模板,自动生成正式合同文本",
    "调用行业协会平台接口,自动生成合同备案表",
    "支持模板维护(可由甲方业务侧自助上传 / 修改)",
]
y = Inches(2.4)
for t in left_bullets:
    add_textbox(s, Inches(1.0), y, Inches(0.4), Inches(0.5),
                "▸", size=20, bold=True, color=ACCENT)
    add_textbox(s, Inches(1.4), y + Inches(0.05), Inches(4.9), Inches(0.9),
                t, size=14, color=GREY_DARK)
    y += Inches(0.95)

# 右卡片
add_rounded(s, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.3), WHITE, line=SUCCESS)
add_rect(s, Inches(6.8), Inches(1.5), Inches(5.8), Inches(0.6), fill=SUCCESS)
add_textbox(s, Inches(7.0), Inches(1.55), Inches(5.4), Inches(0.5),
            "③ 电子流转单与审批", size=18, bold=True, color=WHITE,
            anchor=MSO_ANCHOR.MIDDLE)
right_bullets = [
    "合同生成后自动起单,按部门流转",
    "各部门询问配合费、核对预估总价",
    "支持电子签证,替代纸质签字",
    "审批日志可追溯,符合内控要求",
]
y = Inches(2.4)
for t in right_bullets:
    add_textbox(s, Inches(7.1), y, Inches(0.4), Inches(0.5),
                "▸", size=20, bold=True, color=SUCCESS)
    add_textbox(s, Inches(7.5), y + Inches(0.05), Inches(4.9), Inches(0.9),
                t, size=14, color=GREY_DARK)
    y += Inches(0.95)
add_footer(s, 8, TOTAL)

# ===================== 9. 核心功能 4 + 5 BPM 同步 / AI 分析 =====================
s = prs.slides.add_slide(blank)
add_page_title(s, "五 · 核心功能 ④ ⑤ · BPM 数据同步与 AI 分析")

# 左:BPM
add_rounded(s, Inches(0.7), Inches(1.5), Inches(5.8), Inches(5.3), WHITE, line=PRIMARY)
add_rect(s, Inches(0.7), Inches(1.5), Inches(5.8), Inches(0.6), fill=PRIMARY)
add_textbox(s, Inches(0.9), Inches(1.55), Inches(5.4), Inches(0.5),
            "④ BPM 系统对接", size=18, bold=True, color=WHITE,
            anchor=MSO_ANCHOR.MIDDLE)
bpm = [
    "BPM 为单一系统,同时开放内网 / 外网登录",
    "覆盖 BPM 商机 + BPM 销售合同两个模块,不再二次录入",
    "需对接 BPM 提供的 API 或数据库写入方式",
    "需与 BPM 厂商共同确认字段映射 + 接口认证方式",
]
y = Inches(2.4)
for t in bpm:
    add_textbox(s, Inches(1.0), y, Inches(0.4), Inches(0.5),
                "▸", size=20, bold=True, color=PRIMARY)
    add_textbox(s, Inches(1.4), y + Inches(0.05), Inches(4.9), Inches(0.9),
                t, size=14, color=GREY_DARK)
    y += Inches(0.95)

# 右:AI
add_rounded(s, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.3), WHITE, line=WARN)
add_rect(s, Inches(6.8), Inches(1.5), Inches(5.8), Inches(0.6), fill=WARN)
add_textbox(s, Inches(7.0), Inches(1.55), Inches(5.4), Inches(0.5),
            "⑤ AI 月度数据分析", size=18, bold=True, color=WHITE,
            anchor=MSO_ANCHOR.MIDDLE)
ai = [
    "每月自动汇总合同执行情况、价格 / 折扣分布",
    "识别异常合同(价格异常、配合费异常)",
    "形成图表化分析报告,推送管理层",
    "数据维度可由业务侧灵活配置",
]
y = Inches(2.4)
for t in ai:
    add_textbox(s, Inches(7.1), y, Inches(0.4), Inches(0.5),
                "▸", size=20, bold=True, color=WARN)
    add_textbox(s, Inches(7.5), y + Inches(0.05), Inches(4.9), Inches(0.9),
                t, size=14, color=GREY_DARK)
    y += Inches(0.95)
add_footer(s, 9, TOTAL)

# ===================== 10. 老方法流程图 =====================
s = prs.slides.add_slide(blank)
add_page_title(s, "六 · ① 老方法流程", "现状 — 同一份信息,人工反复录入跨多系统")

# 流程节点(老方法)
old_steps = [
    ("01\n获取\n许可证", "纸质 / 截图\n/ PDF", "人工"),
    ("02\n抄录到\n协会网", "外网 · 老旧\n字段繁多", "人工"),
    ("03\n抄录到\n管理平台", "内网 · 老旧\n字段重复", "人工"),
    ("04\n填写\nWord 合同", "套模板\n手敲价格", "人工"),
    ("05\n打印 +\n纸质流转", "盖章 + 签字\n+ 走廊跑动", "人工"),
    ("06\n抄录到\nBPM 商机", "BPM 模块一\n商机录入", "人工"),
    ("07\n抄录到\nBPM 销售", "BPM 模块二\n销售合同录入", "人工"),
    ("08\n抄录到\n合同归档", "管理平台\n再次录入", "人工"),
]
n = len(old_steps)
step_w = Inches(1.45)
step_h = Inches(2.4)
total_w = step_w * n + Inches(0.05) * (n - 1)
start_x = (Inches(13.333) - total_w) / 2
y = Inches(1.6)
for i, (title, desc, role) in enumerate(old_steps):
    x = start_x + i * (step_w + Inches(0.05))
    add_rounded(s, x, y, step_w, step_h, WHITE, line=WARN)
    add_rect(s, x, y, step_w, Inches(1.0), fill=WARN)
    add_textbox(s, x, y + Inches(0.08), step_w, Inches(0.9),
                title, size=12, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(s, x + Inches(0.05), y + Inches(1.1), step_w - Inches(0.1),
                Inches(0.7), desc, size=10, color=GREY_DARK,
                align=PP_ALIGN.CENTER)
    # 人工标签
    tag_w = step_w - Inches(0.2)
    add_rounded(s, x + Inches(0.1), y + Inches(1.85), tag_w, Inches(0.4),
                WARN)
    add_textbox(s, x + Inches(0.1), y + Inches(1.87), tag_w, Inches(0.35),
                role, size=11, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if i < n - 1:
        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                 x + step_w - Inches(0.02), y + Inches(0.35),
                                 Inches(0.12), Inches(0.3))
        arr.fill.solid()
        arr.fill.fore_color.rgb = WARN
        arr.line.fill.background()

# 痛点小结
add_rect(s, Inches(0.7), Inches(4.55), Inches(11.9), Inches(2.45), fill=GREY_LIGHT)
add_textbox(s, Inches(0.9), Inches(4.7), Inches(11.6), Inches(0.5),
            "为什么慢、为什么错?", size=18, bold=True, color=WARN)
pain_pts = [
    "8 个环节全部依赖人工录入,重复抄写同一组字段",
    "跨内网(管理平台)、外网(协会网)、纸质(流转单)三种介质",
    "BPM 商机 + BPM 销售合同 两个模块都要再录一次",
    "纸质流转影响速度,签字签证无审计日志",
    "字段口径靠人记忆,出错后追溯困难",
]
y = Inches(5.2)
for t in pain_pts:
    add_textbox(s, Inches(0.95), y, Inches(0.4), Inches(0.4),
                "✕", size=14, bold=True, color=WARN)
    add_textbox(s, Inches(1.3), y, Inches(11.2), Inches(0.4),
                t, size=12, color=GREY_DARK)
    y += Inches(0.35)
add_footer(s, 10, TOTAL)

# ===================== 11. 新平台流程图 =====================
s = prs.slides.add_slide(blank)
add_page_title(s, "六 · ② 新平台流程", "上线后 — 人工只录一次,系统自动联动各下游")

new_steps = [
    ("01\n上传\n许可证", "图片 / 截图\n/ PDF", "人工"),
    ("02\nOCR\n识别", "自动抓取\n字段", "系统"),
    ("03\n人工\n补全", "价格 / 折扣 /\n配合费", "人工"),
    ("04\n自动生成\n合同", "Word +\n协会备案表", "系统"),
    ("05\n电子\n流转审批", "各部门 +\n电子签证", "系统"),
    ("06\n对接\nBPM", "商机 + 销售\n合同自动写入", "系统"),
    ("07\n归档 +\nAI 分析", "数据沉淀 +\n月度报告", "系统"),
]
n = len(new_steps)
step_w = Inches(1.65)
step_h = Inches(2.4)
total_w = step_w * n + Inches(0.1) * (n - 1)
start_x = (Inches(13.333) - total_w) / 2
y = Inches(1.6)
for i, (title, desc, role) in enumerate(new_steps):
    x = start_x + i * (step_w + Inches(0.1))
    is_manual = (role == "人工")
    border = PRIMARY if is_manual else SUCCESS
    header = PRIMARY if is_manual else SUCCESS
    add_rounded(s, x, y, step_w, step_h, WHITE, line=border)
    add_rect(s, x, y, step_w, Inches(1.0), fill=header)
    add_textbox(s, x, y + Inches(0.08), step_w, Inches(0.9),
                title, size=12, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(s, x + Inches(0.05), y + Inches(1.1), step_w - Inches(0.1),
                Inches(0.7), desc, size=10, color=GREY_DARK,
                align=PP_ALIGN.CENTER)
    tag_w = step_w - Inches(0.2)
    tag_color = PRIMARY if is_manual else SUCCESS
    add_rounded(s, x + Inches(0.1), y + Inches(1.85), tag_w, Inches(0.4),
                tag_color)
    add_textbox(s, x + Inches(0.1), y + Inches(1.87), tag_w, Inches(0.35),
                role, size=11, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if i < n - 1:
        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                 x + step_w - Inches(0.02), y + Inches(0.35),
                                 Inches(0.14), Inches(0.3))
        arr.fill.solid()
        arr.fill.fore_color.rgb = SUCCESS
        arr.line.fill.background()

# 改进对比
add_rect(s, Inches(0.7), Inches(4.55), Inches(11.9), Inches(2.45), fill=GREY_LIGHT)
add_textbox(s, Inches(0.9), Inches(4.7), Inches(11.6), Inches(0.5),
            "相比老方法的改进", size=18, bold=True, color=SUCCESS)
gain_pts = [
    "8 个环节 → 7 个环节;其中仅 2 步保留人工(上传 + 价格补全)",
    "全程电子化,纸质流转单 / 签字本被电子流转 + 电子签证取代",
    "BPM 商机 + 销售合同两个模块由平台一次写入",
    "字段口径由系统统一,数据天然结构化,可被 AI 月度分析",
    "新平台可二次开发,后续业务变化通过插件 / 脚本快速跟进",
]
y = Inches(5.2)
for t in gain_pts:
    add_textbox(s, Inches(0.95), y, Inches(0.4), Inches(0.4),
                "✓", size=14, bold=True, color=SUCCESS)
    add_textbox(s, Inches(1.3), y, Inches(11.2), Inches(0.4),
                t, size=12, color=GREY_DARK)
    y += Inches(0.35)
add_footer(s, 11, TOTAL)

# ===================== 12. 参考 · 已有线上版合同流转单 =====================
s = prs.slides.add_slide(blank)
add_page_title(s, "参考 · 已有线上版合同流转单系统",
               "此前已完成开发,新平台在此基础上扩展 OCR + 多端同步能力")

screenshots = [
    ("登录页", r"C:\Users\admin\cc_test\contract_flow\s1_login.png",
     "多角色登录(项目承接人/负责人/部门/营销/技术质量/管理员)"),
    ("列表页", r"C:\Users\admin\cc_test\contract_flow\s2_list.png",
     "按状态筛选(草稿→审批中→已完成),支持导出 Excel / 打印 PDF"),
    ("表单编辑页", r"C:\Users\admin\cc_test\contract_flow\s3_form.png",
     "完整电子表单+合同文件 OCR 识别+附件上传+电子签名审批"),
]

col_w = Inches(4.0)
col_x = [Inches(0.5), Inches(4.7), Inches(9.1)]
col_y = Inches(1.3)
col_h = Inches(5.9)

for i, (title, path, desc) in enumerate(screenshots):
    x = col_x[i]
    # 标题条
    add_rect(s, x, col_y, col_w, Inches(0.4), fill=ACCENT)
    add_textbox(s, x, col_y + Inches(0.02), col_w, Inches(0.36),
                title, size=14, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 图片(按栏宽自适应)
    pic = s.shapes.add_picture(path, x, col_y + Inches(0.5),
                                width=col_w)
    pic_h = pic.height
    # 说明文字(图片下方)
    txt_y = col_y + Inches(0.5) + pic_h + Inches(0.15)
    txt_h = col_h - Inches(0.5) - pic_h - Inches(0.15)
    add_textbox(s, x, txt_y, col_w, txt_h,
                desc, size=12, color=GREY_DARK,
                align=PP_ALIGN.CENTER)

add_footer(s, 12, TOTAL)

# ===================== 13. 网络拓扑与跨网数据流 =====================
s = prs.slides.add_slide(blank)
add_page_title(s, "七、网络拓扑与跨网数据流", "新平台部署在内网,需安全打通外网协会 + 内/外可达 BPM")

# 内网区(左)
inner_x, inner_y = Inches(0.7), Inches(1.4)
inner_w, inner_h = Inches(5.6), Inches(5.0)
add_rounded(s, inner_x, inner_y, inner_w, inner_h, RGBColor(0xEA, 0xF1, 0xFA),
            line=PRIMARY)
add_rect(s, inner_x, inner_y, inner_w, Inches(0.5), fill=PRIMARY)
add_textbox(s, inner_x, inner_y + Inches(0.05), inner_w, Inches(0.4),
            "内网 · 公司局域网", size=16, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# 内网组件
inner_items = [
    ("新合同管理平台", "主系统(本期建设)", SUCCESS, True),
    ("老管理平台", "逐步退役", WARN, False),
    ("BPM 系统", "内 / 外网均可登录", PRIMARY, False),
    ("OCR 服务", "可选自研 / 内网部署", ACCENT, False),
]
ix = inner_x + Inches(0.3)
iy = inner_y + Inches(0.8)
for name, desc, color, is_main in inner_items:
    h = Inches(0.95) if is_main else Inches(0.85)
    add_rounded(s, ix, iy, inner_w - Inches(0.6), h, WHITE, line=color)
    add_rect(s, ix, iy, Inches(0.1), h, fill=color)
    add_textbox(s, ix + Inches(0.25), iy + Inches(0.1),
                inner_w - Inches(0.9), Inches(0.4),
                name, size=14, bold=True, color=color)
    add_textbox(s, ix + Inches(0.25), iy + Inches(0.45),
                inner_w - Inches(0.9), Inches(0.4),
                desc, size=11, color=GREY_MED)
    iy += h + Inches(0.15)

# 外网区(右)
outer_x = Inches(7.0)
outer_y = inner_y
outer_w, outer_h = inner_w, inner_h
add_rounded(s, outer_x, outer_y, outer_w, outer_h, RGBColor(0xFD, 0xEF, 0xEC),
            line=WARN)
add_rect(s, outer_x, outer_y, outer_w, Inches(0.5), fill=WARN)
add_textbox(s, outer_x, outer_y + Inches(0.05), outer_w, Inches(0.4),
            "外网 · 公网", size=16, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

outer_items = [
    ("行业协会平台", "外网 · 老旧 · 合同备案表", WARN, False),
    ("住建委系统", "施工许可证查询 / 截图来源", ACCENT, False),
    ("BPM 外网入口", "员工外出时通过外网访问", PRIMARY, False),
]
ix = outer_x + Inches(0.3)
iy = outer_y + Inches(0.8)
for name, desc, color, _ in outer_items:
    h = Inches(0.9)
    add_rounded(s, ix, iy, outer_w - Inches(0.6), h, WHITE, line=color)
    add_rect(s, ix, iy, Inches(0.1), h, fill=color)
    add_textbox(s, ix + Inches(0.25), iy + Inches(0.1),
                outer_w - Inches(0.9), Inches(0.4),
                name, size=14, bold=True, color=color)
    add_textbox(s, ix + Inches(0.25), iy + Inches(0.45),
                outer_w - Inches(0.9), Inches(0.4),
                desc, size=11, color=GREY_MED)
    iy += h + Inches(0.15)

# 中间跨网通道
gate_x = Inches(6.3)
gate_y = Inches(3.4)
gate_w = Inches(0.7)
gate_h = Inches(1.0)
add_rounded(s, gate_x, gate_y, gate_w, gate_h, ACCENT)
add_textbox(s, gate_x, gate_y + Inches(0.05), gate_w, gate_h - Inches(0.1),
            "跨网\n数据\n通道", size=10, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# 双向箭头
for dy, dirn in [(Inches(0.2), "right"), (Inches(0.5), "left")]:
    shape_type = MSO_SHAPE.RIGHT_ARROW if dirn == "right" else MSO_SHAPE.LEFT_ARROW
    if dirn == "right":
        arr = s.shapes.add_shape(shape_type,
                                 gate_x + gate_w + Inches(0.02),
                                 gate_y + dy,
                                 Inches(0.6), Inches(0.2))
    else:
        arr = s.shapes.add_shape(shape_type,
                                 gate_x - Inches(0.62),
                                 gate_y + dy,
                                 Inches(0.6), Inches(0.2))
    arr.fill.solid()
    arr.fill.fore_color.rgb = ACCENT
    arr.line.fill.background()

# 底部说明
add_rect(s, Inches(0.7), Inches(6.55), Inches(11.9), Inches(0.5), fill=GREY_LIGHT)
add_textbox(s, Inches(0.9), Inches(6.6), Inches(11.6), Inches(0.45),
            "跨网通道方案待与开发方共同设计:专线 / API 网关 / 文件交换 / RPA;需满足公司安全合规要求",
            size=12, bold=True, color=PRIMARY, anchor=MSO_ANCHOR.MIDDLE)
add_footer(s, 13, TOTAL)

# ===================== 14. 待讨论 / 待澄清问题 =====================
s = prs.slides.add_slide(blank)
add_page_title(s, "八、与开发方待讨论问题", "希望软件公司在初步沟通中给出方向性意见")

topics = [
    ("跨网数据交换方案", "内网新平台 → 外网协会网的安全打通方式;是否走专线 / 网关 / 文件交换"),
    ("协会网对接方式", "协会网老旧、只有外网;有无 API,还是只能 Web 自动化 / RPA"),
    ("BPM 系统对接", "BPM 是否提供 API / 数据库写入;字段映射 + 接口认证方式"),
    ("OCR 实现方式", "自研 / 第三方 OCR / 多模态大模型;准确率与成本权衡"),
    ("老管理平台数据迁移", "历史合同数据是否搬运,搬运范围 / 字段口径如何对齐"),
    ("Word 模板机制", "模板由谁维护,变量如何标记,业务侧能否自助修改"),
    ("权限与组织架构", "审批流配置方式,多角色 / 多部门支持"),
    ("部署与运维", "内网部署形式,运维由甲方还是开发方承担"),
    ("二次开发能力", "提供 API / 插件 / Webhook 哪种粒度,后续业务变化如何跟"),
    ("交付节奏", "MVP 范围、里程碑、试点上线时间"),
]
cols = 2
card_w = Inches(5.9)
card_h = Inches(1.05)
start_x = Inches(0.7)
start_y = Inches(1.4)
gap_x = Inches(0.2)
gap_y = Inches(0.1)
for i, (t, d) in enumerate(topics):
    r, c = divmod(i, cols)
    x = start_x + c * (card_w + gap_x)
    y = start_y + r * (card_h + gap_y)
    is_critical = (i < 3)  # 前 3 项是网络/对接相关,标红
    border = WARN if is_critical else ACCENT
    add_rounded(s, x, y, card_w, card_h, WHITE, line=border)
    add_rect(s, x, y, Inches(0.1), card_h, fill=border)
    add_textbox(s, x + Inches(0.25), y + Inches(0.1), card_w - Inches(0.4),
                Inches(0.4), t, size=14, bold=True,
                color=WARN if is_critical else PRIMARY)
    add_textbox(s, x + Inches(0.25), y + Inches(0.5), card_w - Inches(0.4),
                card_h - Inches(0.55), d, size=11, color=GREY_DARK)
add_footer(s, 13, TOTAL)

# ===================== 14. 下一步 =====================
s = prs.slides.add_slide(blank)
add_page_title(s, "九、后续推进", "建议的下一步动作")

steps = [
    ("第 1 步", "本次沟通", "确认理解一致,识别认知差异;约定下次会议议题"),
    ("第 2 步", "需求细化", "由甲方业务侧补齐字段清单、模板样例、现有系统截图"),
    ("第 3 步", "技术方案", "由开发方给出整体技术方案 / 工期估算 / 报价"),
    ("第 4 步", "MVP 范围", "共同界定 MVP(建议优先:OCR + 合同生成 + 流转单)"),
    ("第 5 步", "合作启动", "签约,试点上线,迭代扩展"),
]
y = Inches(1.6)
for i, (no, t, d) in enumerate(steps):
    add_rounded(s, Inches(0.7), y, Inches(1.5), Inches(0.9), PRIMARY)
    add_textbox(s, Inches(0.7), y + Inches(0.2), Inches(1.5), Inches(0.5),
                no, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rounded(s, Inches(2.3), y, Inches(10.3), Inches(0.9), WHITE, line=PRIMARY)
    add_textbox(s, Inches(2.5), y + Inches(0.1), Inches(3), Inches(0.4),
                t, size=16, bold=True, color=PRIMARY)
    add_textbox(s, Inches(2.5), y + Inches(0.45), Inches(9.8), Inches(0.4),
                d, size=13, color=GREY_DARK)
    y += Inches(1.05)
add_footer(s, 15, TOTAL)

# ===================== 16. 附录① 协会平台 + 合同文本制作 =====================
s = prs.slides.add_slide(blank)
add_page_title(s, "附录① · 当前录入界面",
               "行业协会平台 / 合同文本制作 — 现状截图")
col_w = Inches(6.15)
col_h = Inches(5.9)
col_y = Inches(1.3)
add_module_block(s, Inches(0.5), col_y, col_w, col_h,
                 "行业协会平台 — 录入 + 勾选参数",
                 ["image1.png", "image2.png"], label_color=WARN)
add_module_block(s, Inches(6.85), col_y, col_w, col_h,
                 "合同文本制作 — Word 模板填写",
                 ["image3.png", "image4.png"], label_color=ACCENT)
add_footer(s, 15, TOTAL)

# ===================== 16. 附录② 管理平台 + BPM 商机 =====================
s = prs.slides.add_slide(blank)
add_page_title(s, "附录② · 当前录入界面",
               "管理平台 / BPM 商机 — 现状截图")
add_module_block(s, Inches(0.5), col_y, col_w, col_h,
                 "管理平台 — 录入 + 合同归档(老旧 · 待替代)",
                 ["image5.png", "image6.png", "image7.png"], label_color=WARN)
add_module_block(s, Inches(6.85), col_y, col_w, col_h,
                 "BPM 商机 — 商机数据录入",
                 ["image8.png", "image9.png"], label_color=PRIMARY)
add_footer(s, 17, TOTAL)

# ===================== 18. 附录③ BPM 销售合同 =====================
s = prs.slides.add_slide(blank)
add_page_title(s, "附录③ · 当前录入界面",
               "BPM 销售合同 — 现状截图")
# 4 张图,2x2 网格
quad_w = Inches(6.15)
quad_h = Inches(2.85)
quad_y = Inches(1.3)
def add_image_card(slide, x, y, w, h, title, file_name):
    add_rect(slide, x, y, w, Inches(0.4), fill=PRIMARY)
    add_textbox(slide, x + Inches(0.15), y + Inches(0.02), w - Inches(0.3),
                Inches(0.36), title, size=12, bold=True, color=WHITE,
                anchor=MSO_ANCHOR.MIDDLE)
    box_y = y + Inches(0.45)
    box_h = h - Inches(0.45)
    add_rect(slide, x, box_y, w, box_h, fill=WHITE, line=GREY_MED, line_w=0.5)
    fit_picture(slide, img(file_name), x + Inches(0.05), box_y + Inches(0.05),
                w - Inches(0.1), box_h - Inches(0.1))

sales_imgs = [
    ("销售合同 · 基础信息", "image10.png"),
    ("销售合同 · 主要字段", "image11.png"),
    ("销售合同 · 费用明细", "image12.png"),
    ("销售合同 · 流程节点", "image13.png"),
]
for i, (t, f) in enumerate(sales_imgs):
    r, c = divmod(i, 2)
    x = Inches(0.5) + c * (quad_w + Inches(0.2))
    y = quad_y + r * (quad_h + Inches(0.2))
    add_image_card(s, x, y, quad_w, quad_h, t, f)
add_footer(s, 17, TOTAL)

# ===================== 18. 施工许可证 — OCR 信息源(压轴页) =====================
s = prs.slides.add_slide(blank)
add_page_title(s, "施工许可证 — OCR 信息源",
               "新平台所有自动化流程的起点 — 一张图,八个字段")

# 左侧大图
left_x = Inches(0.5)
left_y = Inches(1.3)
left_w = Inches(7.5)
left_h = Inches(5.9)
add_rect(s, left_x, left_y, left_w, left_h, fill=WHITE, line=SUCCESS, line_w=1.5)
fit_picture(s, img("image14.jpeg"),
            left_x + Inches(0.1), left_y + Inches(0.1),
            left_w - Inches(0.2), left_h - Inches(0.2))

# 右侧字段标签
right_x = Inches(8.3)
right_y = left_y
right_w = Inches(4.55)
add_rect(s, right_x, right_y, right_w, Inches(0.5), fill=SUCCESS)
add_textbox(s, right_x, right_y + Inches(0.05), right_w, Inches(0.4),
            "OCR 自动抓取的 8 个字段", size=14, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
ocr_fields = [
    "报建号", "施工许可证编号",
    "工程名称", "建设地址",
    "建设单位", "施工单位",
    "设计单位", "监理单位",
]
fy = right_y + Inches(0.7)
for i, f in enumerate(ocr_fields):
    r, c = divmod(i, 2)
    fx = right_x + c * Inches(2.3)
    fyi = fy + r * Inches(0.65)
    add_rounded(s, fx, fyi, Inches(2.2), Inches(0.55), WHITE, line=SUCCESS)
    add_textbox(s, fx, fyi + Inches(0.13), Inches(2.2), Inches(0.4),
                f, size=12, bold=True, color=SUCCESS,
                align=PP_ALIGN.CENTER)

# 底部小结
sum_y = fy + 4 * Inches(0.65) + Inches(0.2)
add_rect(s, right_x, sum_y, right_w, Inches(1.4), fill=GREY_LIGHT)
add_textbox(s, right_x + Inches(0.15), sum_y + Inches(0.15),
            right_w - Inches(0.3), Inches(0.4),
            "为什么这页放最后?", size=12, bold=True, color=PRIMARY)
add_textbox(s, right_x + Inches(0.15), sum_y + Inches(0.55),
            right_w - Inches(0.3), Inches(0.8),
            "8 个字段被 OCR 一次读出,后续 6+ 平台 / 模板 / 流转单 / BPM,"
            "全部由系统自动复用 — 这就是新平台代替老管理平台的价值原点。",
            size=11, color=GREY_DARK)
add_footer(s, 19, TOTAL)

# ---------- 保存 ----------
out = r"C:\Users\admin\Desktop\scan\新平台开发需求.pptx"
prs.save(out)
print(f"PPT 已生成: {out}")
print(f"共 {TOTAL} 页")
