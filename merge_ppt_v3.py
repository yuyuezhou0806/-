"""
Merge two PowerPoint presentations by properly copying all slide parts
and their related parts (layouts, masters, themes, media).
"""
from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from lxml import etree
from copy import deepcopy
import os


def copy_part_recursive(source_part, target_pkg, part_map=None):
    """Copy a part and all its related parts recursively."""
    if part_map is None:
        part_map = {}

    if source_part in part_map:
        return part_map[source_part]

    # Create new part with same content type
    new_part = target_pkg.part_factory(
        source_part.content_type,
        source_part.blob
    )
    part_map[source_part] = new_part

    # Copy relationships
    for rel in source_part.rels.values():
        target_rel_part = copy_part_recursive(rel.target_part, target_pkg, part_map)
        new_part.relate_to(target_rel_part, rel.reltype)

    return new_part


def merge_presentations_v3(output_path, pres_paths):
    """Merge presentations by copying slide parts with all related parts."""

    # Create new presentation
    merged = Presentation()

    # We need to use the internal package to add parts
    merged_pkg = merged.part.package

    for pres_path in pres_paths:
        prs = Presentation(pres_path)
        source_pkg = prs.part.package

        # Map to track copied parts
        part_map = {}

        for slide in prs.slides:
            slide_part = slide.part

            # Copy slide part and all its related parts
            new_slide_part = copy_part_recursive(slide_part, merged_pkg, part_map)

            # Add slide to presentation
            # Get the slide layout part from the new slide part's relationships
            layout_part = None
            for rel in new_slide_part.rels.values():
                if rel.reltype == RT.SLIDE_LAYOUT:
                    layout_part = rel.target_part
                    break

            if layout_part is None:
                print(f"Warning: No layout found for slide, skipping")
                continue

            # Add slide to presentation using the presentation part
            presentation_part = merged.part
            rId = presentation_part.relate_to(new_slide_part, RT.SLIDE)

            # Add sldId to sldIdLst
            sldIdLst = presentation_part._element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}sldIdLst')
            if sldIdLst is not None:
                nsmap = sldIdLst.nsmap
                p_ns = nsmap.get(None, 'http://schemas.openxmlformats.org/presentationml/2006/main')
                r_ns = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

                # Find max id
                max_id = 256
                for sldId in sldIdLst.findall(f'{{{p_ns}}}sldId'):
                    id_val = int(sldId.get('id', 256))
                    max_id = max(max_id, id_val)

                new_sldId = etree.SubElement(sldIdLst, f'{{{p_ns}}}sldId')
                new_sldId.set('id', str(max_id + 1))
                new_sldId.set(f'{{{r_ns}}}id', rId)

    merged.save(output_path)
    print(f"Saved: {output_path}")
    print(f"Total slides: {len(merged.slides)}")


# Alternative: use python-pptx's internal slide duplication
from pptx.oxml.ns import qn
from pptx.opc.package import XmlPart


def duplicate_slide(pres, index):
    """Duplicate a slide within the same presentation."""
    source = pres.slides[index]

    # Get blank layout
    try:
        blank_layout = pres.slide_layouts[6]
    except IndexError:
        blank_layout = pres.slide_layouts[-1]

    blank = pres.slides.add_slide(blank_layout)
    blank_shapes = blank.shapes

    # Copy all shapes from source
    for shape in source.shapes:
        el = shape.element
        new_el = deepcopy(el)
        blank_shapes._spTree.insert_element_before(new_el, 'p:extLst')

    return blank


def merge_by_rebuild(output_path, *input_paths):
    """
    Merge by reading each presentation, extracting images,
    and rebuilding slides with images in a new presentation.
    """
    merged = Presentation()
    merged.slide_width = Inches(13.333)
    merged.slide_height = Inches(7.5)

    for input_path in input_paths:
        prs = Presentation(input_path)
        for slide in prs.slides:
            # Add blank slide
            try:
                blank_layout = merged.slide_layouts[6]
            except:
                blank_layout = merged.slide_layouts[-1]
            new_slide = merged.slides.add_slide(blank_layout)

            # Copy each shape
            for shape in slide.shapes:
                el = shape.element
                new_el = deepcopy(el)
                new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')

    merged.save(output_path)
    return merged


# Let's try the simplest approach: use python-pptx to load and save,
# but properly handle images by using the part system.
from pptx.util import Inches


def merge_ppts(output, *inputs):
    """Merge presentations preserving images using part-level operations."""
    from pptx.opc.package import PartFactory

    merged = Presentation()
    merged.slide_width = Inches(13.333)
    merged.slide_height = Inches(7.5)
    merged_pkg = merged.part.package

    for input_path in inputs:
        prs = Presentation(input_path)

        for slide in prs.slides:
            slide_part = slide.part

            # Find the slide layout relationship
            layout_rel = None
            for rel in slide_part.rels.values():
                if rel.reltype == RT.SLIDE_LAYOUT:
                    layout_rel = rel
                    break

            if layout_rel is None:
                continue

            layout_part = layout_rel.target_part

            # Check if layout part already exists in merged
            merged_layout = None
            for existing_layout in merged.slide_layouts:
                if existing_layout.part.blob == layout_part.blob:
                    merged_layout = existing_layout
                    break

            # If not found, we need to add it - but python-pptx doesn't support
            # adding arbitrary layouts. Use blank layout instead.
            if merged_layout is None:
                try:
                    merged_layout = merged.slide_layouts[6]
                except:
                    merged_layout = merged.slide_layouts[-1]

            # Create new slide
            new_slide = merged.slides.add_slide(merged_layout)

            # Copy shape elements
            for shape in slide.shapes:
                el = shape.element
                new_el = deepcopy(el)
                new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')

    merged.save(output)
    return merged


# Actually the issue is that deepcopy doesn't preserve image references.
# Let's use a completely different approach: direct zip manipulation
# with proper part copying.

import zipfile
import tempfile
import shutil


def copy_part_content(source_pkg, target_pkg, source_part, copied_parts=None):
    """Copy a part and all its related parts from source to target package."""
    if copied_parts is None:
        copied_parts = {}

    if source_part in copied_parts:
        return copied_parts[source_part]

    # Create new part in target package
    new_part = target_pkg.part_factory(
        source_part.content_type,
        source_part.blob
    )
    copied_parts[source_part] = new_part

    # Recursively copy related parts
    for rel in source_part.rels.values():
        target_related = copy_part_content(source_pkg, target_pkg, rel.target_part, copied_parts)
        new_part.relate_to(target_related, rel.reltype)

    return new_part


def merge_presentations_final(output_path, *input_paths):
    """Robust merge using python-pptx's part system."""
    merged = Presentation()
    merged.slide_width = Inches(13.333)
    merged.slide_height = Inches(7.5)

    for input_path in input_paths:
        prs = Presentation(input_path)
        source_pkg = prs.part.package
        target_pkg = merged.part.package

        copied_parts = {}

        for slide in prs.slides:
            slide_part = slide.part

            # Copy the slide part and all related parts
            new_slide_part = copy_part_content(source_pkg, target_pkg, slide_part, copied_parts)

            # Add slide to presentation
            pres_part = merged.part
            rId = pres_part.relate_to(new_slide_part, RT.SLIDE)

            # Update sldIdLst
            sldIdLst_elem = pres_part._element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}sldIdLst')
            if sldIdLst_elem is not None:
                p_ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
                r_ns = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

                max_id = 256
                for sldId in sldIdLst_elem.findall(f'{{{p_ns}}}sldId'):
                    id_val = int(sldId.get('id', 256))
                    max_id = max(max_id, id_val)

                new_sldId = etree.SubElement(sldIdLst_elem, f'{{{p_ns}}}sldId')
                new_sldId.set('id', str(max_id + 1))
                new_sldId.set(f'{{{r_ns}}}id', rId)

    merged.save(output_path)
    return merged


# Run it
output = r"c:\Users\admin\cc_test\系统操作流程.pptx"
input1 = r"c:\Users\admin\cc_test\商机制作流程.pptx"
input2 = r"c:\Users\admin\cc_test\客户申领流程.pptx"

merged = merge_presentations_final(output, input1, input2)
print(f"Merged: {len(merged.slides)} slides")

# Verify
import os
size = os.path.getsize(output)
print(f"File size: {size / 1024:.0f} KB")

# Test loading
verify = Presentation(output)
print(f"Verification: {len(verify.slides)} slides loaded successfully")

for i, slide in enumerate(verify.slides):
    pic_count = len([s for s in slide.shapes if s.shape_type == 13])
    if pic_count > 0:
        print(f"  Slide {i+1}: {pic_count} images")
