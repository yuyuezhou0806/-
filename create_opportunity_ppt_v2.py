from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

IMG_DIR = r"c:\Users\admin\cc_test"

def add_top_bar(slide, color=RGBColor(0x1a, 0x5f, 0x7a)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

def add_title(slide, text, top=Inches(0.25), color=RGBColor(0x1a, 0x5f, 0x7a)):
    tb = slide.shapes.add_textbox(Inches(0.4), top, Inches(12.5), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = color
    return tb

def add_bullet_text(slide, items, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_before = Pt(12)
        p.space_after = Pt(6)
    return tb

def add_image_slide(slide, img_path, img_left, img_top, max_width, max_height):
    try:
        with Image.open(img_path) as img:
            orig_w, orig_h = img.size
        aspect = orig_w / orig_h
        if max_width / max_height > aspect:
            h = max_height
            w = h * aspect
        else:
            w = max_width
            h = w / aspect
        pic = slide.shapes.add_picture(img_path, img_left, img_top, width=w, height=h)
        return pic
    except Exception as e:
        print(f"Error adding image {img_path}: {e}")
        return None

# ========== 第1页：封面 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = RGBColor(0x1a, 0x5f, 0x7a)
bg.line.fill.background()
tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.333), Inches(1.5))
p = tb.text_frame.paragraphs[0]
p.text = "商机制作流程"
p.font.size = Pt(54)
p.font.bold = True
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p.alignment = PP_ALIGN.CENTER
p2 = tb.text_frame.add_paragraph()
p2.text = "系统操作手册"
p2.font.size = Pt(28)
p2.font.color.rgb = RGBColor(0xCC, 0xE5, 0xFF)
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(20)

# ========== 第2页：字段类型说明 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_top_bar(slide)
add_title(slide, "字段类型说明")
add_bullet_text(slide, [
    "🔴 红色边框 — 输入框：手动输入文字或数字",
    "      商机名称、预计成交金额、投资额或面积",
    "",
    "🟡 黄色边框 — 选择框：下拉/弹窗选择",
    "      子公司、业务经理、业态、区域、日期等",
    "",
    "🟢 绿色边框 — 特殊字段：系统关联字段",
    "      商机重要程度、客户名称（后续详细介绍）"
], Inches(0.4), Inches(1.1), Inches(5.5), Inches(5))
add_image_slide(slide, f"{IMG_DIR}\\39dd04aab749169b695f5e317492332d.png",
                Inches(6.2), Inches(1.0), Inches(6.8), Inches(6.2))

# ========== 第3页：红色输入框 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_top_bar(slide, RGBColor(0xD4, 0x4B, 0x4B))
add_title(slide, "第一步：填写输入框（红色边框）", color=RGBColor(0xD4, 0x4B, 0x4B))
add_bullet_text(slide, [
    "商机名称：填写项目全称",
    "      例：嘉松南路城桂路交叉口新建工程",
    "",
    "预计成交金额(元)：填写预估金额",
    "      例：10,000.00",
    "",
    "投资额或面积：填写投资规模",
    "      无具体数据时可填 /"
], Inches(0.4), Inches(1.1), Inches(5.5), Inches(5))
add_image_slide(slide, f"{IMG_DIR}\\39dd04aab749169b695f5e317492332d.png",
                Inches(6.2), Inches(1.0), Inches(6.8), Inches(6.2))

# ========== 第4页：黄色选择框-基本信息 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_top_bar(slide, RGBColor(0xC7, 0x94, 0x00))
add_title(slide, "第二步：选择基本信息（黄色边框）", color=RGBColor(0xC7, 0x94, 0x00))
add_bullet_text(slide, [
    "子公司：选择承接主体  例：中测行",
    "业务经理：选择负责人  例：卫天勇",
    "业态：选择业务类型  例：工程检测",
    "一级服务名录：选择工程类别  例：房屋建筑工程",
    "产品名录：选择检测类型  例：建筑材料检测",
    "项目大区：选择地理大区  例：华东地区",
    "项目区域：选择省市  例：上海"
], Inches(0.4), Inches(1.1), Inches(5.5), Inches(5.5))
add_image_slide(slide, f"{IMG_DIR}\\39dd04aab749169b695f5e317492332d.png",
                Inches(6.2), Inches(1.0), Inches(6.8), Inches(6.2))

# ========== 第5页：黄色选择框-日期与其他 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_top_bar(slide, RGBColor(0xC7, 0x94, 0x00))
add_title(slide, "第三步：选择日期与其他字段（黄色边框）", color=RGBColor(0xC7, 0x94, 0x00))
add_bullet_text(slide, [
    "📅 开标日期：项目公开招标日期",
    "      例：2026-05-19",
    "📅 预计定标日期：预计确定中标方日期",
    "      例：2026-05-24",
    "📅 预计签约日期：预计签订合同日期",
    "      例：2026-05-26",
    "",
    "🎯 中标可能性：评估中标概率  例：中"
], Inches(0.4), Inches(1.1), Inches(5.5), Inches(5.5))
add_image_slide(slide, f"{IMG_DIR}\\39dd04aab749169b695f5e317492332d.png",
                Inches(6.2), Inches(1.0), Inches(6.8), Inches(6.2))

# ========== 第6页：客户信息与参与人员 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_top_bar(slide)
add_title(slide, "第四步：填写客户信息与参与人员")
add_bullet_text(slide, [
    "客户名称：从系统客户库中选择",
    "      例：上海廷瑞建设工程有限公司",
    "      （绿色特殊框，关联CRM客户主数据）",
    "",
    "商机阶段：当前所处销售阶段",
    "      例：需求分析(20%)",
    "",
    "参与人员：添加项目成员",
    "      例：卫天勇 / 中测行经营室 / 100% / 是"
], Inches(0.4), Inches(1.1), Inches(5.5), Inches(5.5))
add_image_slide(slide, f"{IMG_DIR}\\d1cd2336e65a528736e42f061cd3f8c2.png",
                Inches(6.2), Inches(1.0), Inches(6.8), Inches(6.2))

# ========== 第7页：特殊字段（绿色）==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_top_bar(slide, RGBColor(0x28, 0xA7, 0x45))
add_title(slide, "特殊字段详解（绿色边框）", color=RGBColor(0x28, 0xA7, 0x45))
add_bullet_text(slide, [
    "🟢 商机重要程度",
    "      系统根据成交金额、客户等级等自动判定",
    "      例：一般商机",
    "      影响审批流程层级和资源支持力度",
    "",
    "🟢 客户名称",
    "      必须从系统客户库中选择，非自由文本",
    "      确保项目与客户档案准确关联",
    "      影响后续合同、回款、客户统计等业务"
], Inches(0.4), Inches(1.1), Inches(5.5), Inches(5.5))
add_image_slide(slide, f"{IMG_DIR}\\d1cd2336e65a528736e42f061cd3f8c2.png",
                Inches(6.2), Inches(1.0), Inches(6.8), Inches(6.2))

# ========== 第8页：提交前检查清单 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_top_bar(slide)
add_title(slide, "提交前检查清单")
add_bullet_text(slide, [
    "□ 商机名称是否完整准确？",
    "□ 子公司、业务经理是否已选择？",
    "□ 业态、服务名录、产品名录是否匹配？",
    "□ 项目大区/区域是否正确？",
    "□ 预计成交金额是否填写？",
    "□ 开标/定标/签约三个日期是否已选择？",
    "□ 客户名称是否已从系统中选择？",
    "□ 参与人员及提成分配是否配置？",
    "□ 主要负责人是否已指定？",
    "",
    "✅ 所有必填项（带 * 号）完成后，点击「提交」"
], Inches(0.4), Inches(1.1), Inches(12.5), Inches(5.5))

# 保存
output = r"c:\Users\admin\cc_test\商机制作流程.pptx"
prs.save(output)
print(f"PPT已生成：{output}")
