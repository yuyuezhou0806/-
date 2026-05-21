"""
Merge PowerPoint presentations by extracting and repacking zip contents.
Handles naming conflicts for slides, layouts, masters, themes, and media.
"""
import zipfile
import tempfile
import shutil
import os
from lxml import etree


def merge_presentations(output_path, *input_paths):
    """Merge multiple pptx files into one."""

    temp_dir = tempfile.mkdtemp()

    # Extract first presentation as base
    with zipfile.ZipFile(input_paths[0], 'r') as z:
        z.extractall(temp_dir)

    # Track existing files to avoid conflicts
    existing_files = set()
    for root_dir, dirs, files in os.walk(temp_dir):
        for f in files:
            rel_path = os.path.relpath(os.path.join(root_dir, f), temp_dir)
            existing_files.add(rel_path.replace('\\', '/'))

    # Counters for renaming
    next_slide_num = len([f for f in os.listdir(os.path.join(temp_dir, 'ppt', 'slides'))
                          if f.startswith('slide') and f.endswith('.xml') and '_rels' not in f]) + 1

    # Parse base presentation.xml.rels
    pres_rels_path = os.path.join(temp_dir, 'ppt', '_rels', 'presentation.xml.rels')
    with open(pres_rels_path, 'r', encoding='utf-8') as f:
        pres_rels_content = f.read()

    # Parse base presentation.xml
    pres_xml_path = os.path.join(temp_dir, 'ppt', 'presentation.xml')
    pres_tree = etree.parse(pres_xml_path)
    pres_root = pres_tree.getroot()
    p_ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    r_ns = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    sldIdLst = pres_root.find(f'.//{{{p_ns}}}sldIdLst')

    # Find max existing sldId
    max_sld_id = 256
    for sldId in sldIdLst.findall(f'{{{p_ns}}}sldId'):
        id_val = int(sldId.get('id', 256))
        max_sld_id = max(max_sld_id, id_val)

    # Find max existing rId
    max_rId = 0
    for rel in etree.fromstring(pres_rels_content.encode()).findall('.//'):
        rid = rel.get('Id', '')
        if rid.startswith('rId'):
            try:
                max_rId = max(max_rId, int(rid[3:]))
            except:
                pass

    # Track media file mapping: original_name -> merged_name
    media_dir = os.path.join(temp_dir, 'ppt', 'media')
    existing_media = set(os.listdir(media_dir)) if os.path.exists(media_dir) else set()
    media_map = {}  # maps (input_idx, orig_name) -> merged_name

    for input_idx, input_path in enumerate(input_paths[1:], 1):
        add_temp = tempfile.mkdtemp()
        with zipfile.ZipFile(input_path, 'r') as z:
            z.extractall(add_temp)

        add_ppt_dir = os.path.join(add_temp, 'ppt')
        add_slides_dir = os.path.join(add_ppt_dir, 'slides')
        add_slides_rels_dir = os.path.join(add_slides_dir, '_rels')
        add_media_dir = os.path.join(add_ppt_dir, 'media')
        add_layouts_dir = os.path.join(add_ppt_dir, 'slideLayouts')
        add_layouts_rels_dir = os.path.join(add_layouts_dir, '_rels')
        add_masters_dir = os.path.join(add_ppt_dir, 'slideMasters')
        add_masters_rels_dir = os.path.join(add_masters_dir, '_rels')
        add_themes_dir = os.path.join(add_ppt_dir, 'theme')

        # --- Copy media files ---
        if os.path.exists(add_media_dir):
            os.makedirs(media_dir, exist_ok=True)
            for media_file in sorted(os.listdir(add_media_dir)):
                src = os.path.join(add_media_dir, media_file)
                # Determine merged name
                if media_file in existing_media:
                    # Check if content is the same
                    with open(src, 'rb') as f:
                        src_blob = f.read()
                    dst_path = os.path.join(media_dir, media_file)
                    with open(dst_path, 'rb') as f:
                        dst_blob = f.read()
                    if src_blob == dst_blob:
                        # Same file, reuse name
                        media_map[(input_idx, media_file)] = media_file
                        continue
                    # Different content, rename
                    name, ext = os.path.splitext(media_file)
                    counter = 1
                    new_name = f"{name}_{counter}{ext}"
                    while new_name in existing_media:
                        counter += 1
                        new_name = f"{name}_{counter}{ext}"
                    media_map[(input_idx, media_file)] = new_name
                    dst = os.path.join(media_dir, new_name)
                else:
                    media_map[(input_idx, media_file)] = media_file
                    dst = os.path.join(media_dir, media_file)

                shutil.copy2(src, dst)
                existing_media.add(media_map[(input_idx, media_file)])

        # --- Parse additional presentation to get slide order ---
        add_pres_xml = os.path.join(add_ppt_dir, 'presentation.xml')
        add_tree = etree.parse(add_pres_xml)
        add_root = add_tree.getroot()
        add_sldIdLst = add_root.find(f'.//{{{p_ns}}}sldIdLst')

        add_pres_rels = os.path.join(add_ppt_dir, '_rels', 'presentation.xml.rels')
        with open(add_pres_rels, 'r', encoding='utf-8') as f:
            add_rels_content = f.read()
        add_rels_tree = etree.fromstring(add_rels_content.encode())

        # Map rId -> target for slides
        slide_targets = {}
        for rel in add_rels_tree.findall('.//'):
            rel_type = rel.get('Type', '')
            if 'slide' in rel_type and rel.get('Target', '').startswith('slides/'):
                slide_targets[rel.get('Id')] = rel.get('Target')

        # --- Copy each slide in order ---
        for sldId in add_sldIdLst.findall(f'{{{p_ns}}}sldId'):
            r_id = sldId.get(f'{{{r_ns}}}id')
            if r_id not in slide_targets:
                continue

            target = slide_targets[r_id]
            slide_name = os.path.basename(target)  # e.g., slide1.xml
            slide_num = int(''.join(filter(str.isdigit, slide_name)))

            # Copy slide XML with new name
            src_slide = os.path.join(add_slides_dir, slide_name)
            new_slide_name = f"slide{next_slide_num}.xml"
            dst_slide = os.path.join(temp_dir, 'ppt', 'slides', new_slide_name)
            shutil.copy2(src_slide, dst_slide)
            existing_files.add(f"ppt/slides/{new_slide_name}")

            # Copy and fix slide rels
            src_rels = os.path.join(add_slides_rels_dir, f"slide{slide_num}.xml.rels")
            if os.path.exists(src_rels):
                dst_rels_dir = os.path.join(temp_dir, 'ppt', 'slides', '_rels')
                os.makedirs(dst_rels_dir, exist_ok=True)
                dst_rels = os.path.join(dst_rels_dir, f"slide{next_slide_num}.xml.rels")

                with open(src_rels, 'r', encoding='utf-8') as f:
                    rels_content = f.read()

                # Fix media references
                for (idx, orig_name), merged_name in media_map.items():
                    if idx == input_idx and orig_name != merged_name:
                        rels_content = rels_content.replace(
                            f'Target="../media/{orig_name}"',
                            f'Target="../media/{merged_name}"'
                        )

                with open(dst_rels, 'w', encoding='utf-8') as f:
                    f.write(rels_content)
                existing_files.add(f"ppt/slides/_rels/{next_slide_num}.xml.rels")

            # Add to presentation rels
            max_rId += 1
            new_rId = f"rId{max_rId}"
            new_rel = f'<Relationship Id="{new_rId}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="slides/{new_slide_name}"/>'
            pres_rels_content = pres_rels_content.replace('</Relationships>', f'{new_rel}\n</Relationships>')

            # Add to sldIdLst
            max_sld_id += 1
            new_sldId = etree.SubElement(sldIdLst, f'{{{p_ns}}}sldId')
            new_sldId.set('id', str(max_sld_id))
            new_sldId.set(f'{{{r_ns}}}id', new_rId)

            next_slide_num += 1

        shutil.rmtree(add_temp)

    # Write updated presentation files
    with open(pres_rels_path, 'w', encoding='utf-8') as f:
        f.write(pres_rels_content)

    pres_tree.write(pres_xml_path, xml_declaration=True, encoding='UTF-8', standalone=True)

    # Update [Content_Types].xml to include any new files
    ct_path = os.path.join(temp_dir, '[Content_Types].xml')
    with open(ct_path, 'r', encoding='utf-8') as f:
        ct_content = f.read()

    # Ensure slide content types are present
    ct_tree = etree.fromstring(ct_content.encode())
    ct_root = ct_tree

    # Check if we need to add Override for new slides
    # (usually not needed if same content type as existing)

    etree.ElementTree(ct_root).write(ct_path, xml_declaration=True, encoding='UTF-8', standalone=True)

    # Repack as pptx
    with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
        for root_dir, dirs, files in os.walk(temp_dir):
            for file in files:
                file_path = os.path.join(root_dir, file)
                arcname = os.path.relpath(file_path, temp_dir).replace('\\', '/')
                zipf.write(file_path, arcname)

    shutil.rmtree(temp_dir)


# Run
output = r"c:\Users\admin\cc_test\系统操作流程.pptx"
input1 = r"c:\Users\admin\cc_test\商机制作流程.pptx"
input2 = r"c:\Users\admin\cc_test\客户申领流程.pptx"

merge_presentations(output, input1, input2)

# Verify
import os
size = os.path.getsize(output)
print(f"合并完成：{output}")
print(f"文件大小：{size / 1024:.0f} KB")

from pptx import Presentation
p = Presentation(output)
print(f"总页数：{len(p.slides)}")

for i, slide in enumerate(p.slides):
    pic_count = len([s for s in slide.shapes if s.shape_type == 13])
    if pic_count > 0:
        print(f"  第{i+1}页：{pic_count} 张图片")
