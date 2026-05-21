"""
Merge PowerPoint presentations - fixed version with Content_Types.xml update.
"""
import zipfile
import tempfile
import shutil
import os
from lxml import etree


def merge_presentations(output_path, *input_paths):
    temp_dir = tempfile.mkdtemp()

    # Extract first presentation as base
    with zipfile.ZipFile(input_paths[0], 'r') as z:
        z.extractall(temp_dir)

    # Track media
    media_dir = os.path.join(temp_dir, 'ppt', 'media')
    media_hashes = {}
    if os.path.exists(media_dir):
        for f in os.listdir(media_dir):
            with open(os.path.join(media_dir, f), 'rb') as fp:
                media_hashes[f] = fp.read()

    slides_dir = os.path.join(temp_dir, 'ppt', 'slides')
    next_slide_num = len([f for f in os.listdir(slides_dir)
                          if f.startswith('slide') and f.endswith('.xml')]) + 1

    # Parse presentation.xml
    pres_xml_path = os.path.join(temp_dir, 'ppt', 'presentation.xml')
    pres_tree = etree.parse(pres_xml_path)
    pres_root = pres_tree.getroot()
    p_ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    r_ns = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    sldIdLst = pres_root.find(f'.//{{{p_ns}}}sldIdLst')

    max_sld_id = 256
    for sldId in sldIdLst.findall(f'{{{p_ns}}}sldId'):
        max_sld_id = max(max_sld_id, int(sldId.get('id', 256)))

    # Parse presentation.xml.rels
    pres_rels_path = os.path.join(temp_dir, 'ppt', '_rels', 'presentation.xml.rels')
    pres_rels_tree = etree.parse(pres_rels_path)
    pres_rels_root = pres_rels_tree.getroot()
    rel_ns = 'http://schemas.openxmlformats.org/package/2006/relationships'

    max_rId = 0
    for rel in pres_rels_root.findall(f'{{{rel_ns}}}Relationship'):
        rid = rel.get('Id', '')
        if rid.startswith('rId'):
            try:
                max_rId = max(max_rId, int(rid[3:]))
            except:
                pass

    # Content types to add
    new_content_types = []
    slide_ct = 'application/vnd.openxmlformats-officedocument.presentationml.slide+xml'

    media_map = {}

    for input_idx, input_path in enumerate(input_paths[1:], 1):
        add_temp = tempfile.mkdtemp()
        with zipfile.ZipFile(input_path, 'r') as z:
            z.extractall(add_temp)

        add_ppt = os.path.join(add_temp, 'ppt')
        add_slides = os.path.join(add_ppt, 'slides')
        add_slides_rels = os.path.join(add_slides, '_rels')
        add_media = os.path.join(add_ppt, 'media')

        # Copy media
        if os.path.exists(add_media):
            os.makedirs(media_dir, exist_ok=True)
            for media_file in sorted(os.listdir(add_media)):
                src = os.path.join(add_media, media_file)
                with open(src, 'rb') as f:
                    blob = f.read()

                found = False
                for existing_name, existing_blob in media_hashes.items():
                    if blob == existing_blob:
                        media_map[(input_idx, media_file)] = existing_name
                        found = True
                        break

                if not found:
                    dst_name = media_file
                    counter = 1
                    while dst_name in media_hashes:
                        name, ext = os.path.splitext(media_file)
                        dst_name = f"{name}_{counter}{ext}"
                        counter += 1
                    dst = os.path.join(media_dir, dst_name)
                    shutil.copy2(src, dst)
                    media_hashes[dst_name] = blob
                    media_map[(input_idx, media_file)] = dst_name

        # Parse additional presentation
        add_pres = os.path.join(add_ppt, 'presentation.xml')
        add_tree = etree.parse(add_pres)
        add_root = add_tree.getroot()
        add_sldIdLst = add_root.find(f'.//{{{p_ns}}}sldIdLst')

        add_rels_path = os.path.join(add_ppt, '_rels', 'presentation.xml.rels')
        add_rels_tree = etree.parse(add_rels_path)
        add_rels_root = add_rels_tree.getroot()

        slide_targets = {}
        for rel in add_rels_root.findall(f'{{{rel_ns}}}Relationship'):
            rel_type = rel.get('Type', '')
            target = rel.get('Target', '')
            if 'slide' in rel_type and target.startswith('slides/'):
                slide_targets[rel.get('Id')] = target

        for sldId in add_sldIdLst.findall(f'{{{p_ns}}}sldId'):
            r_id = sldId.get(f'{{{r_ns}}}id')
            if r_id not in slide_targets:
                continue

            target = slide_targets[r_id]
            slide_name = os.path.basename(target)
            slide_num = int(''.join(filter(str.isdigit, slide_name)))

            # Copy slide XML
            src_slide = os.path.join(add_slides, slide_name)
            new_name = f"slide{next_slide_num}.xml"
            dst_slide = os.path.join(slides_dir, new_name)
            shutil.copy2(src_slide, dst_slide)

            # Track content type
            new_content_types.append(f'/ppt/slides/{new_name}')

            # Copy and fix slide rels
            src_rels = os.path.join(add_slides_rels, f"slide{slide_num}.xml.rels")
            if os.path.exists(src_rels):
                dst_rels_dir = os.path.join(temp_dir, 'ppt', 'slides', '_rels')
                os.makedirs(dst_rels_dir, exist_ok=True)
                dst_rels = os.path.join(dst_rels_dir, f"slide{next_slide_num}.xml.rels")

                with open(src_rels, 'r', encoding='utf-8') as f:
                    rels_content = f.read()

                for (idx, orig), merged in media_map.items():
                    if idx == input_idx and orig != merged:
                        rels_content = rels_content.replace(
                            f'Target="../media/{orig}"',
                            f'Target="../media/{merged}"'
                        )

                with open(dst_rels, 'w', encoding='utf-8') as f:
                    f.write(rels_content)

            # Add to presentation rels
            max_rId += 1
            new_rId = f"rId{max_rId}"
            new_rel = etree.SubElement(pres_rels_root, f'{{{rel_ns}}}Relationship')
            new_rel.set('Id', new_rId)
            new_rel.set('Type', 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide')
            new_rel.set('Target', f'slides/{new_name}')

            # Add to sldIdLst
            max_sld_id += 1
            new_sldId = etree.SubElement(sldIdLst, f'{{{p_ns}}}sldId')
            new_sldId.set('id', str(max_sld_id))
            new_sldId.set(f'{{{r_ns}}}id', new_rId)

            next_slide_num += 1

        shutil.rmtree(add_temp)

    # Write presentation files
    pres_rels_tree.write(pres_rels_path, xml_declaration=True, encoding='UTF-8', standalone=True)
    pres_tree.write(pres_xml_path, xml_declaration=True, encoding='UTF-8', standalone=True)

    # Update [Content_Types].xml
    ct_path = os.path.join(temp_dir, '[Content_Types].xml')
    ct_tree = etree.parse(ct_path)
    ct_root = ct_tree.getroot()
    ct_ns = 'http://schemas.openxmlformats.org/package/2006/content-types'

    for part_path in new_content_types:
        override = etree.SubElement(ct_root, f'{{{ct_ns}}}Override')
        override.set('PartName', part_path)
        override.set('ContentType', slide_ct)

    ct_tree.write(ct_path, xml_declaration=True, encoding='UTF-8', standalone=True)

    # Repack
    with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
        for root_dir, dirs, files in os.walk(temp_dir):
            for file in files:
                file_path = os.path.join(root_dir, file)
                arcname = os.path.relpath(file_path, temp_dir).replace('\\', '/')
                zipf.write(file_path, arcname)

    shutil.rmtree(temp_dir)


# Run
output = r"c:\Users\admin\cc_test\系统操作流程.pptx"
input1 = r"c:\Users\admin\cc_test\商机制作流程.pptx"
input2 = r"c:\Users\admin\cc_test\客户申领流程.pptx"

merge_presentations(output, input1, input2)

# Verify
import os
size = os.path.getsize(output)
print(f"合并完成：{output}")
print(f"文件大小：{size / 1024:.0f} KB")

from pptx import Presentation
p = Presentation(output)
print(f"总页数：{len(p.slides)}")

for i, slide in enumerate(p.slides):
    pic_count = len([s for s in slide.shapes if s.shape_type == 13])
    if pic_count > 0:
        print(f"  第{i+1}页：{pic_count} 张图片")
