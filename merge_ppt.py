from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from copy import deepcopy
import lxml

def copy_slide(source_prs, target_prs, index):
    """Copy a slide from source presentation to target presentation"""
    source_slide = source_prs.slides[index]

    # Get the slide layout from source
    slide_layout = source_slide.slide_layout

    # Try to find matching layout in target, or use blank
    try:
        target_layout = target_prs.slide_layouts[6]  # blank layout
    except:
        target_layout = target_prs.slide_layouts[-1]

    # Add new slide with blank layout
    new_slide = target_prs.slides.add_slide(target_layout)

    # Copy all shapes from source slide
    for shape in source_slide.shapes:
        el = shape.element
        new_el = deepcopy(el)
        new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')

    return new_slide

# Create new presentation
merged_prs = Presentation()
merged_prs.slide_width = Inches(13.333)
merged_prs.slide_height = Inches(7.5)

# Load both presentations
opp_prs = Presentation(r"c:\Users\admin\cc_test\商机制作流程.pptx")
cust_prs = Presentation(r"c:\Users\admin\cc_test\客户申领流程.pptx")

print(f"商机PPT: {len(opp_prs.slides)} 页")
print(f"客户PPT: {len(cust_prs.slides)} 页")

# Copy all slides from 商机制作流程
for i in range(len(opp_prs.slides)):
    copy_slide(opp_prs, merged_prs, i)

# Add a divider slide
slide = merged_prs.slides.add_slide(merged_prs.slide_layouts[6])
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), merged_prs.slide_width, merged_prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = RGBColor(0x1a, 0x5f, 0x7a)
bg.line.fill.background()
tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.333), Inches(1.5))
p = tb.text_frame.paragraphs[0]
p.text = "第二部分"
p.font.size = Pt(54)
p.font.bold = True
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p.alignment = PP_ALIGN.CENTER
p2 = tb.text_frame.add_paragraph()
p2.text = "客户申领流程"
p2.font.size = Pt(32)
p2.font.color.rgb = RGBColor(0xCC, 0xE5, 0xFF)
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(20)

# Copy all slides from 客户申领流程
for i in range(len(cust_prs.slides)):
    copy_slide(cust_prs, merged_prs, i)

# Save
output = r"c:\Users\admin\cc_test\系统操作流程.pptx"
merged_prs.save(output)
print(f"合并完成：{output}")
print(f"总计：{len(merged_prs.slides)} 页")
