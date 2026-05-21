from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle=""):
    """添加标题页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)

    # 背景色
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(0x1a, 0x5f, 0x7a)
    background.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(28)
        p2.font.color.rgb = RGBColor(0xCC, 0xE5, 0xFF)
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(20)

    return slide

def add_content_slide(prs, title, content_items, notes=None, highlight_colors=None):
    """添加内容页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 顶部色条
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.15)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0x1a, 0x5f, 0x7a)
    bar.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1a, 0x5f, 0x7a)

    # 内容区域
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.333), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = item
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_before = Pt(16)
        p.space_after = Pt(8)
        p.level = 0

    # 添加备注
    if notes:
        notes_slide = slide.notes_slide
        notes_text_frame = notes_slide.notes_text_frame
        notes_text_frame.text = notes

    return slide

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """添加双栏内容页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 顶部色条
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.15)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0x1a, 0x5f, 0x7a)
    bar.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1a, 0x5f, 0x7a)

    # 左侧标题
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(5.8), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xD4, 0x4B, 0x4B)  # 红色 - 输入框

    # 左侧内容
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(5.8), Inches(5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "▸ " + item
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_before = Pt(12)

    # 右侧标题
    right_title_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.3), Inches(5.8), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xC7, 0x94, 0x00)  # 金色 - 选择框

    # 右侧内容
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.9), Inches(5.8), Inches(5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "▸ " + item
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_before = Pt(12)

    return slide

def add_special_slide(prs, title, items):
    """添加特殊字段说明页（绿色主题）"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 顶部绿色色条
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.15)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0x28, 0xA7, 0x45)
    bar.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x28, 0xA7, 0x45)

    # 内容
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.333), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_before = Pt(18)
        p.space_after = Pt(8)

    return slide

# ========== 开始构建PPT ==========

# 第1页：封面
add_title_slide(prs, "商机制作流程指南", "系统操作手册 · 2026-05-19")

# 第2页：流程总览
add_content_slide(prs, "流程总览", [
    "1. 填写商机基本信息 — 商机名称、金额、日期等核心字段",
    "2. 选择关联选项 — 子公司、业务经理、业态、区域等下拉选择",
    "3. 填写客户信息 — 客户名称、联系人、联系方式",
    "4. 配置参与人员 — 添加项目成员、分配提成比例",
    "5. 确认特殊字段 — 商机重要程度、客户名称等系统关联字段",
    "6. 提交审批 — 检查必填项后提交进入审批流程"
])

# 第3页：字段类型说明
add_content_slide(prs, "字段类型说明", [
    "",
    "🔴 红色边框 — 输入框：需要手动输入文字或数字",
    "      例如：商机名称、预计成交金额、投资额或面积",
    "",
    "🟡 黄色边框 — 选择框：点击后从下拉列表或弹窗中选择",
    "      例如：子公司、业务经理、业态、项目区域、日期",
    "",
    "🟢 绿色边框 — 特殊字段：系统关联或带有特殊业务逻辑的字段",
    "      例如：商机重要程度、客户名称（后续详细介绍）"
])

# 第4页：步骤一 - 输入框字段（红色）
add_two_column_slide(prs,
    "步骤一：填写输入框字段（红色边框）",
    "必填输入项",
    [
        "商机名称：填写项目全称",
        "      例：嘉松南路城桂路交叉口新建工程",
        "",
        "预计成交金额(元)：填写预估金额",
        "      例：10,000.00",
        "",
        "投资额或面积：填写投资规模",
        "      例：/（如无可填斜杠）"
    ],
    "填写要点",
    [
        "商机名称建议包含：",
        "  · 道路/地点名称",
        "  · 工程类型关键词",
        "",
        "金额注意单位为元，",
        "支持自动千分位格式",
        "",
        "投资额无具体数据时",
        "可用 '/' 占位"
    ]
)

# 第5页：步骤二 - 选择框字段（黄色）- 基本信息
add_two_column_slide(prs,
    "步骤二：选择基本信息字段（黄色边框）",
    "组织与人员",
    [
        "子公司：选择承接主体",
        "      例：中测行",
        "",
        "业务经理：选择负责人",
        "      例：卫天勇",
        "      （选择后自动带出所属部门）",
        "",
        "业态：选择业务类型",
        "      例：工程检测"
    ],
    "服务与区域",
    [
        "一级服务名录：选择工程类别",
        "      例：房屋建筑工程",
        "",
        "产品名录：选择检测类型",
        "      例：建筑材料检测",
        "",
        "项目大区：选择地理大区",
        "      例：华东地区",
        "",
        "项目区域：选择省市",
        "      例：上海"
    ]
)

# 第6页：步骤三 - 日期选择（黄色）
add_content_slide(prs, "步骤三：选择日期字段（黄色边框）", [
    "",
    "📅 开标日期：项目公开招标的日期",
    "      例：2026-05-19",
    "",
    "📅 预计定标日期：预计确定中标方的日期",
    "      例：2026-05-24",
    "",
    "📅 预计签约日期：预计签订合同的日期",
    "      例：2026-05-26",
    "",
    "💡 提示：日期通过点击输入框弹出日历选择，也可手动输入格式 YYYY-MM-DD"
])

# 第7页：步骤四 - 其他选择字段（黄色）
add_content_slide(prs, "步骤四：其他选择字段（黄色边框）", [
    "",
    "🎯 中标可能性：评估中标概率",
    "      选项：高 / 中 / 低",
    "      例：中",
    "",
    "📊 商机阶段：当前所处销售阶段",
    "      例：需求分析(20%)",
    "      （通常由系统根据流程自动更新）",
    "",
    "👤 参与人员 → 人员姓名：",
    "      从组织架构中选择项目成员",
    "      例：卫天勇（自动带出部门：中测行经营室）"
])

# 第8页：步骤五 - 客户信息
add_two_column_slide(prs,
    "步骤五：填写客户信息",
    "客户名称（绿色边框 — 特殊字段）",
    [
        "客户名称：选择或搜索客户",
        "      例：上海廷瑞建设工程有限公司",
        "",
        "✅ 该字段为绿色特殊框，",
        "说明它关联客户主数据系统",
        "",
        "选择后会自动关联该客户的",
        "历史项目、信用评级等信息"
    ],
    "联系信息",
    [
        "客户联系人：",
        "      选择该客户下的联系人",
        "      （可在右侧 + 添加新联系人）",
        "",
        "客户电话：",
        "      自动带出联系人电话",
        "      或手动填写",
        "",
        "💡 如客户不在系统中，",
        "需先在客户管理模块新建"
    ]
)

# 第9页：步骤六 - 参与人员配置
add_content_slide(prs, "步骤六：配置参与人员与提成", [
    "",
    "👥 人员姓名：选择项目参与人员（黄色选择框）",
    "      例：卫天勇",
    "",
    "🏢 人员部门：自动带出所选人员所属部门",
    "      例：中测行经营室",
    "",
    "💰 提成分配比例：填写该成员的提成占比",
    "      例：100.00%（多人时总和应为100%）",
    "",
    "⭐ 是否主要负责人：标记项目主负责人",
    "      例：是（单选，只有一个主要负责人）"
])

# 第10页：特殊字段详解（绿色）
add_special_slide(prs, "特殊字段详解（绿色边框）", [
    "",
    "🟢 商机重要程度",
    "      系统根据成交金额、客户等级、战略价值等自动判定",
    "      例：一般商机（绿色标签显示）",
    "      其他级别可能包括：重点商机、战略商机等",
    "      该字段影响审批流程层级和资源支持力度",
    "",
    "🟢 客户名称",
    "      关联CRM客户主数据，非自由文本输入",
    "      必须从系统客户库中选择",
    "      确保项目与客户档案准确关联",
    "      影响后续合同、回款、客户统计等业务",
    "",
    "💡 绿色字段通常由系统规则驱动，不建议随意修改"
])

# 第11页：完整字段对照表
add_content_slide(prs, "完整字段对照速查表", [
    "🔴 输入框（手动输入）：商机名称、预计成交金额(元)、投资额或面积",
    "",
    "🟡 选择框（下拉/弹窗选择）：子公司、业务经理、业态、一级服务名录、",
    "      产品名录、项目大区、项目区域、中标可能性、开标日期、预计定标日期、",
    "      预计签约日期、人员姓名",
    "",
    "🟢 特殊字段（系统关联）：商机重要程度、客户名称",
    "",
    "⚪ 自动带出（无需填写）：业务经理所属部门、人员部门、商机阶段"
])

# 第12页：提交前检查清单
add_content_slide(prs, "提交前检查清单", [
    "",
    "□ 商机名称是否完整准确？",
    "□ 子公司、业务经理是否已选择？",
    "□ 业态、服务名录、产品名录是否匹配？",
    "□ 项目大区/区域是否正确？",
    "□ 预计成交金额是否填写？",
    "□ 三个关键日期是否已选择？",
    "□ 客户名称是否已从系统中选择？",
    "□ 参与人员及提成分配是否配置？",
    "□ 主要负责人是否已指定？",
    "",
    "✅ 所有必填项（带 * 号）完成后，点击右上角「提交」按钮"
])

# 第13页：结尾
add_title_slide(prs, "谢谢", "如有疑问请联系系统管理员")

# 保存
output_path = r"c:\Users\admin\cc_test\商机制作流程指南.pptx"
prs.save(output_path)
print(f"PPT已生成：{output_path}")
